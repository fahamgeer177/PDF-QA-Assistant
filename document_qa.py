import os
import openai
from langchain.chains import RetrievalQA
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import Chroma 
from langchain_openai import OpenAIEmbeddings, OpenAI  

from docx import Document
from pptx import Presentation


conversation_retrieval_chain = None
chat_history = []
llm = None
embeddings = None

def init_llm():
    global llm, embeddings

    openai.api_key = os.getenv("OPENAI_API_KEY")
    if not openai.api_key:
        raise ValueError("OPENAI_API_KEY environment variable is not set.")

   
    llm = OpenAI(openai_api_key=openai.api_key)
    embeddings = OpenAIEmbeddings(openai_api_key=openai.api_key)

UPLOAD_FOLDER = "./uploads"

def delete_old_documents():
    for filename in os.listdir(UPLOAD_FOLDER):
        file_path = os.path.join(UPLOAD_FOLDER, filename)
        try:
            os.remove(file_path)
        except Exception as e:
            print(f"Error deleting file {file_path}: {e}")

def reset_document_processing():
    global conversation_retrieval_chain, chat_history
    conversation_retrieval_chain = None  
    chat_history = [] 

def process_new_document(document_path):
    
    reset_document_processing()

    process_document(document_path)
    print(f"Document processed: {document_path}")

def process_pdf(document_path):
    loader = PyPDFLoader(document_path)
    return loader.load()

def process_docx(document_path):
    doc = Document(document_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    text = text.encode('utf-8', 'ignore').decode('utf-8')
    return [{"page_content": text}]

def process_pptx(document_path):
    prs = Presentation(document_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return [{"page_content": text}]

def process_document(document_path):
    global conversation_retrieval_chain

    file_extension = document_path.split('.')[-1].lower()

    if file_extension == "pdf":
        documents = process_pdf(document_path)
    elif file_extension == "docx":
        documents = process_docx(document_path)
    elif file_extension == "pptx":
        documents = process_pptx(document_path)
    else:
        raise ValueError(f"Unsupported document format: {file_extension}")

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1024, chunk_overlap=64)
    texts = text_splitter.split_documents(documents)

    db = Chroma.from_documents(texts, embedding=embeddings)

    conversation_retrieval_chain = RetrievalQA.from_chain_type(
        llm=llm,
        chain_type="stuff",
        retriever=db.as_retriever(search_type="mmr", search_kwargs={'k': 6, 'lambda_mult': 0.25}),
        return_source_documents=False,
        input_key="question"
    )

def process_prompt(prompt):
    global conversation_retrieval_chain
    global chat_history

    output = conversation_retrieval_chain.invoke({"question": prompt, "chat_history": chat_history})
    answer = output["result"]

    chat_history.append((prompt, answer))

    return answer

init_llm()

if __name__ == "__main__":
    # For debugging, you can add test calls here or run a simple test.
    print("document_qa.py loaded and ready.")